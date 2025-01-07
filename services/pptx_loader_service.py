from pptx import Presentation
from models import AnnualSummary, QuaterlyMetrics, RevenueByActivity

class PPTXLoaderService:
    def __init__(self, annual_summary_repo, quarterly_metrics_repo, revenue_by_activity_repo, db):
        self.annual_summary_repo = annual_summary_repo
        self.quarterly_metrics_repo = quarterly_metrics_repo
        self.revenue_by_activity_repo = revenue_by_activity_repo
        self.db = db

    def load_pptx_data(self, file_path):
        try:
            presentation = Presentation(file_path)
            self._process_annual_summary(presentation.slides[0])
            self._process_quarterly_metrics(presentation.slides[1])
            self._process_revenue_by_activity(presentation.slides[2])
            
            return True
        except Exception as e:
            self.db.session.rollback()
            raise e

    def _process_annual_summary(self, slide):
        try:
            key_highlights_text = None
            
            # Find the paragraph containing "Key Highlights"
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        print(paragraph.text)
                        if "Key Highlights" in paragraph.text:
                            key_highlights_text = paragraph.text
                            break
                    if key_highlights_text:
                        break
            
            if not key_highlights_text:
                raise ValueError("Key Highlights not found in the slide")

            # Parse the Key Highlights text
            lines = key_highlights_text.split("\n")
            data = {}
            for line in lines:
                if "Total Revenue" in line:
                    data['total_revenue'] = int(line.split(":")[1].strip().replace("$", "").replace(",", ""))
                elif "Total Memberships Sold" in line:
                    data['total_membership_sold'] = int(line.split(":")[1].strip().replace(",", ""))
                elif "Top Location" in line:
                    data['top_location'] = line.split(":")[1].strip()

            if not data:
                raise ValueError("No valid data found in Key Highlights")

            # Map the parsed data to the model
            summary = AnnualSummary(
                total_revenue=data['total_revenue'],
                total_membership_sold=data['total_membership_sold'],
                top_location=data['top_location']
            )
            self.annual_summary_repo.create(summary)
        except Exception as e:
            raise ValueError(f"Error processing annual summary: {str(e)}")


    def _process_quarterly_metrics(self, slide):
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        print(paragraph.text)


            # Extract data from the second slide
            data = self._extract_table_data(slide)
            
            for row in data:
                metrics = QuaterlyMetrics(
                    year=int(row['year']),
                    quarter=row['quarter'],
                    revenue=float(row['revenue']),
                    memberships_sold=int(row['memberships_sold']),
                    duration=int(row['duration'])
                )
                self.quarterly_metrics_repo.create(metrics)
        except Exception as e:
            raise ValueError(f"Error processing quarterly metrics: {str(e)}")

    def _process_revenue_by_activity(self, slide):
        try:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        print(paragraph.text)

                        
            # Extract data from the third slide
            data = self._extract_table_data(slide)
            
            revenue = RevenueByActivity(
                gym=float(data['gym']),
                pool=float(data['pool']),
                tennis_court=float(data['tennis_court']),
                personal_training=float(data['personal_training']),
                others=float(data['others'])
            )
            self.revenue_by_activity_repo.create(revenue)
        except Exception as e:
            raise ValueError(f"Error processing revenue by activity: {str(e)}")

    def _extract_table_data(self, slide):
        """
        Extract data from a slide's table
        Adjust this method based on your actual PowerPoint structure
        """
        for shape in slide.shapes:
            if shape.has_table:
                return self._parse_table(shape.table)
        return {}

    def _parse_table(self, table):
        """
        Parse table data based on the slide type
        This is a placeholder - implement based on actual table structure
        """
        # Implementation depends on actual table structure in the PowerPoint
        pass