from pptx import Presentation
from ..models import db, KeyHighlights, RevenueDistribution

class PPTXLoaderService:
    def __init__(self, db):
        self.db = db

    def load_pptx_data(self, file_path):
        try:
            key_highlights_data, revenue_distribution_data = self._extract_data_from_ppt(file_path)
            self._clear_existing_data()
            
            for record in key_highlights_data:
                highlights = KeyHighlights(
                    total_revenue=record['total_revenue'],
                    membership_sold=record['membership_sold'],
                    top_location=record['top_location']
                )
                self.db.session.add(highlights)

            for record in revenue_distribution_data:
                distribution = RevenueDistribution(
                    gym=record['gym'],
                    pool=record['pool'],
                    tennis_court=record['tennis_court'],
                    personal_training=record['personal_training']
                )
                self.db.session.add(distribution)

            self.db.session.commit()
            return True
        except Exception as e:
            self.db.session.rollback()
            raise e

    def _extract_data_from_ppt(self, file_path):
        key_highlights_data = []
        revenue_distribution_data = []
        
        presentation = Presentation(file_path)
        
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    text = text.split('\n')
                    if "Key Highlights:" in text:
                        total_revenue = int(text[1].split(': $')[1].replace(',', ''))
                        membership_sold = int(text[2].split(': ')[1].replace(',', ''))
                        top_location = text[3].split(': ')[1]
                        key_highlights_data.append({
                            'total_revenue': total_revenue,
                            'membership_sold': membership_sold,
                            'top_location': top_location
                        })
                    elif "Revenue Distribution:" in text:
                        gym = float(text[1].split(': ')[1].replace('%', ''))
                        pool = float(text[2].split(': ')[1].replace('%', ''))
                        tennis_court = float(text[3].split(': ')[1].replace('%', ''))
                        personal_training = float(text[4].split(': ')[1].replace('%', ''))
                        revenue_distribution_data.append({
                            'gym': gym,
                            'pool': pool,
                            'tennis_court': tennis_court,
                            'personal_training': personal_training
                        })

        return key_highlights_data, revenue_distribution_data

    def _clear_existing_data(self):
        try:
            KeyHighlights.query.delete()
            RevenueDistribution.query.delete()
            self.db.session.commit()
        except Exception:
            self.db.session.rollback()
            raise